"""
Generate a stunning UniVerse UMS PowerPoint presentation.
Uses dark theme with purple-cyan gradient accents, minimal & premium look.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
BG_DARK       = RGBColor(0x0B, 0x0B, 0x1E)   # deep navy-black
BG_CARD       = RGBColor(0x14, 0x14, 0x2E)   # card surface
ACCENT_PURPLE = RGBColor(0x7C, 0x6E, 0xF7)   # primary purple
ACCENT_BLUE   = RGBColor(0x5B, 0x9C, 0xF6)   # secondary blue
ACCENT_CYAN   = RGBColor(0x4D, 0xD0, 0xE1)   # cyan accent
ACCENT_LAV    = RGBColor(0xB3, 0x9D, 0xDB)   # lavender
WHITE         = RGBColor(0xF0, 0xF0, 0xFF)
SUBTLE        = RGBColor(0x8A, 0x8A, 0xB0)   # muted text
DIM           = RGBColor(0x5A, 0x5A, 0x80)
CARD_BORDER   = RGBColor(0x2A, 0x2A, 0x4A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────
def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, fill_color, border_color=None, border_w=Pt(0)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.fill.solid()
        shp.line.fill.fore_color.rgb = border_color
        shp.line.width = border_w
    else:
        shp.line.fill.background()
    # round corners
    try:
        shp.adjustments[0] = 0.05
    except Exception:
        pass
    return shp

def add_circle(slide, left, top, size, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_text(slide, left, top, w, h, text, font_size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    try:
        p.line_spacing = Pt(int(font_size * line_spacing))
    except Exception:
        pass
    return txBox

def add_tag(slide, left, top, text, color=ACCENT_CYAN):
    """Small uppercase tag label"""
    tag_w, tag_h = Inches(2.2), Inches(0.35)
    cr = int(int(color) >> 16 & 0xFF)
    cg = int(int(color) >> 8 & 0xFF)
    cb = int(int(color) & 0xFF)
    # tag background pill
    pill = add_rect(slide, left, top, tag_w, tag_h, RGBColor(
        min(cr + 10, 255) // 8, min(cg + 10, 255) // 8, min(cb + 10, 255) // 8
    ), border_color=RGBColor(cr // 3, cg // 3, cb // 3), border_w=Pt(1))
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    add_text(slide, left + Inches(0.15), top + Inches(0.02), tag_w - Inches(0.3), tag_h,
             text.upper(), font_size=8, color=color, bold=True, align=PP_ALIGN.CENTER)

def add_gradient_bar(slide, left, top, width, height=Inches(0.04)):
    """Thin gradient-like accent bar (approx with 3 segments)"""
    seg_w = width // 3
    colors = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN]
    for i, c in enumerate(colors):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + seg_w * i, top, seg_w + Emu(1), height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.fill.background()

def add_decorative_orb(slide, left, top, size, color, opacity_approx=0.15):
    """Adds a translucent decorative circle"""
    r = int(int(color) >> 16 & 0xFF)
    g = int(int(color) >> 8 & 0xFF)
    b = int(int(color) & 0xFF)
    bg_r = int(int(BG_DARK) >> 16 & 0xFF)
    bg_g = int(int(BG_DARK) >> 8 & 0xFF)
    bg_b = int(int(BG_DARK) & 0xFF)
    orb = add_circle(slide, left, top, size, RGBColor(
        int(r * opacity_approx + bg_r * (1 - opacity_approx)),
        int(g * opacity_approx + bg_g * (1 - opacity_approx)),
        int(b * opacity_approx + bg_b * (1 - opacity_approx)),
    ))
    return orb

def add_glass_card(slide, left, top, w, h, icon, title, desc):
    """Reusable glass card with icon, title, description"""
    card = add_rect(slide, left, top, w, h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(1))
    # icon
    add_text(slide, left + Inches(0.25), top + Inches(0.2), Inches(0.5), Inches(0.5),
             icon, font_size=22, align=PP_ALIGN.LEFT)
    # title
    add_text(slide, left + Inches(0.25), top + Inches(0.7), w - Inches(0.5), Inches(0.35),
             title, font_size=13, color=WHITE, bold=True)
    # desc
    add_text(slide, left + Inches(0.25), top + Inches(1.05), w - Inches(0.5), Inches(0.7),
             desc, font_size=10, color=SUBTLE, line_spacing=1.4)

def add_module_card(slide, left, top, w, h, icon, name, sub):
    """Small module card"""
    card = add_rect(slide, left, top, w, h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(1))
    add_text(slide, left, top + Inches(0.15), w, Inches(0.4),
             icon, font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), w, Inches(0.25),
             name, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.78), w, Inches(0.22),
             sub, font_size=7.5, color=DIM, align=PP_ALIGN.CENTER)

def add_feature_row(slide, left, top, dot_color, text, w=Inches(5.5)):
    """Feature list row with dot"""
    row_h = Inches(0.42)
    row_bg = add_rect(slide, left, top, w, row_h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(0.5))
    # dot
    add_circle(slide, left + Inches(0.18), top + Inches(0.14), Inches(0.12), dot_color)
    # text
    add_text(slide, left + Inches(0.42), top + Inches(0.06), w - Inches(0.6), row_h,
             text, font_size=11, color=SUBTLE)
    return row_h + Inches(0.08)

def add_stat_box(slide, left, top, num, label):
    """Stat number + label"""
    add_text(slide, left, top, Inches(1.8), Inches(0.6),
             num, font_size=32, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(slide, left, top + Inches(0.55), Inches(1.8), Inches(0.3),
             label.upper(), font_size=7.5, color=DIM, align=PP_ALIGN.CENTER, bold=True)

def add_tech_badge(slide, left, top, icon, name):
    """Tech stack badge"""
    w, h = Inches(1.65), Inches(0.45)
    badge = add_rect(slide, left, top, w, h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(0.75))
    add_text(slide, left + Inches(0.08), top + Inches(0.05), Inches(0.35), h,
             icon, font_size=13, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.4), top + Inches(0.08), w - Inches(0.5), h,
             name, font_size=9, color=WHITE, bold=False)

def add_arch_box(slide, left, top, icon, title, sub):
    w, h = Inches(1.7), Inches(0.9)
    box = add_rect(slide, left, top, w, h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(1))
    add_text(slide, left, top + Inches(0.08), w, Inches(0.35),
             f"{icon} {title}", font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.45), w, Inches(0.3),
             sub, font_size=7.5, color=DIM, align=PP_ALIGN.CENTER)

def add_arrow(slide, left, top):
    add_text(slide, left, top, Inches(0.35), Inches(0.5),
             "→", font_size=18, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER, bold=True)

def add_pill(slide, left, top, dot_color, text):
    w, h = Inches(1.8), Inches(0.35)
    pill = add_rect(slide, left, top, w, h, BG_CARD, border_color=CARD_BORDER, border_w=Pt(0.5))
    add_circle(slide, left + Inches(0.12), top + Inches(0.1), Inches(0.12), dot_color)
    add_text(slide, left + Inches(0.32), top + Inches(0.04), w - Inches(0.4), h,
             text, font_size=8, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
solid_bg(sl, BG_DARK)

# decorative orbs
add_decorative_orb(sl, Inches(-1.5), Inches(-1), Inches(5), ACCENT_PURPLE, 0.08)
add_decorative_orb(sl, Inches(9), Inches(4), Inches(5), ACCENT_CYAN, 0.06)
add_decorative_orb(sl, Inches(4), Inches(5.5), Inches(3), ACCENT_BLUE, 0.05)

# accent bar top
add_gradient_bar(sl, Inches(4.5), Inches(0), Inches(4.3), Inches(0.04))

# Logo circle
logo_x, logo_y = Inches(5.7), Inches(0.7)
outer = add_circle(sl, logo_x, logo_y, Inches(1.8), RGBColor(0x14, 0x14, 0x35))
outer.line.fill.solid()
outer.line.fill.fore_color.rgb = RGBColor(0x3A, 0x35, 0x70)
outer.line.width = Pt(2)

# Try to add logo image
logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets', 'ums-logo.png')
if os.path.exists(logo_path):
    sl.shapes.add_picture(logo_path, logo_x + Inches(0.3), logo_y + Inches(0.3), Inches(1.2), Inches(1.2))

# Title
add_text(sl, Inches(0), Inches(2.8), SLIDE_W, Inches(0.9),
         "UniVerse", font_size=56, color=ACCENT_PURPLE, bold=True,
         align=PP_ALIGN.CENTER, font_name='Calibri')

# Subtitle
add_text(sl, Inches(0), Inches(3.65), SLIDE_W, Inches(0.5),
         "U N I V E R S I T Y   M A N A G E M E N T   S Y S T E M",
         font_size=13, color=SUBTLE, align=PP_ALIGN.CENTER, font_name='Calibri', bold=True)

# Divider line
add_gradient_bar(sl, Inches(5), Inches(4.3), Inches(3.3), Inches(0.03))

# Description
add_text(sl, Inches(2.5), Inches(4.6), Inches(8.3), Inches(0.7),
         "A next-generation, full-stack platform for managing every aspect\nof academic operations — built with Spring Boot & modern web technologies.",
         font_size=12, color=DIM, align=PP_ALIGN.CENTER, line_spacing=1.6)

# Team names
add_text(sl, Inches(0), Inches(5.8), SLIDE_W, Inches(0.4),
         "Prachi Pandey  ·  Anushka  ·  April 2026",
         font_size=10, color=ACCENT_CYAN, align=PP_ALIGN.CENTER, bold=False)

# Bottom accent
add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(-2), Inches(-1), Inches(4), ACCENT_PURPLE, 0.06)
add_decorative_orb(sl, Inches(10), Inches(5), Inches(4), ACCENT_CYAN, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "The Challenge")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Problem Statement", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.6),
         "Traditional university administration relies on disconnected spreadsheets, paper records, and legacy systems — leading to data silos, human errors, and wasted effort.",
         font_size=12, color=SUBTLE, line_spacing=1.6)

# 4 cards in 2x2 grid
cards = [
    ("📋", "Fragmented Records", "Student, faculty, and financial data scattered across incompatible formats.", ACCENT_PURPLE),
    ("⏱️", "Manual Overhead", "Attendance, grading & fee tracking done by hand — slow and error-prone.", ACCENT_BLUE),
    ("🔒", "No Central Access", "Lack of a single portal for stakeholders to view and manage operations.", ACCENT_CYAN),
    ("📊", "Zero Analytics", "No real-time dashboard or reports for data-driven decision making.", ACCENT_LAV),
]
for i, (icon, title, desc, accent) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.8)
    y = Inches(2.7) + row * Inches(2.2)
    add_glass_card(sl, x, y, Inches(5.4), Inches(1.9), icon, title, desc)
    # Top accent line on card
    bar = slide_shapes = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        x, y, Inches(5.4), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION / INTRODUCING UNIVERSE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(8), Inches(-1), Inches(4), ACCENT_BLUE, 0.06)

add_tag(sl, Inches(0.8), Inches(0.5), "Our Solution")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Introducing UniVerse", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.5),
         "A comprehensive, web-based University Management System that centralizes academics, administration, and finance into one elegant platform.",
         font_size=12, color=SUBTLE, line_spacing=1.6)

add_gradient_bar(sl, Inches(0.8), Inches(2.35), Inches(1.2), Inches(0.03))

features = [
    (ACCENT_PURPLE, "Unified Dashboard — Real-time overview of students, faculty, attendance & financials"),
    (ACCENT_BLUE,   "Role-Based Access — Secure authentication with admin-level controls"),
    (ACCENT_CYAN,   "End-to-End Modules — From admission to graduation, every workflow covered"),
    (ACCENT_LAV,    "Beautiful UI — Glassmorphism design with smooth animations & particle effects"),
    (ACCENT_PURPLE, "RESTful APIs — Clean, well-structured JSON endpoints for all operations"),
    (ACCENT_BLUE,   "Comprehensive Reports — Student, attendance, finance & examination analytics"),
]
y_pos = Inches(2.7)
for dot_color, text in features:
    dy = add_feature_row(sl, Inches(0.8), y_pos, dot_color, text, w=Inches(10))
    y_pos += dy

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — TECH STACK
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(-1), Inches(3), Inches(4), ACCENT_LAV, 0.05)
add_decorative_orb(sl, Inches(10), Inches(-1), Inches(3), ACCENT_CYAN, 0.04)

add_tag(sl, Inches(0.8), Inches(0.5), "Technology")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Tech Stack", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "Built on industry-standard, production-grade technologies.",
         font_size=12, color=SUBTLE)

# Backend section
add_text(sl, Inches(0.8), Inches(2.4), Inches(3), Inches(0.35),
         "⬡  BACKEND", font_size=11, color=ACCENT_CYAN, bold=True)

backend_badges = [("☕", "Java 17"), ("🍃", "Spring Boot 4"), ("🗄️", "Spring Data JPA"),
                  ("🐬", "MySQL 8"), ("🔐", "Spring Security"), ("🧪", "H2 Test DB")]

for i, (icon, name) in enumerate(backend_badges):
    col = i % 3
    row = i // 3
    add_tech_badge(sl, Inches(0.8) + col * Inches(1.8), Inches(2.85) + row * Inches(0.58), icon, name)

# Frontend section
add_text(sl, Inches(6.5), Inches(2.4), Inches(3), Inches(0.35),
         "◆  FRONTEND", font_size=11, color=ACCENT_PURPLE, bold=True)

frontend_badges = [("🌐", "HTML5 / CSS3"), ("⚡", "Vanilla JavaScript"), ("✨", "Particle Effects"),
                   ("🎨", "Glassmorphism"), ("🔤", "Outfit + Inter"), ("📐", "CSS Grid / Flex")]

for i, (icon, name) in enumerate(frontend_badges):
    col = i % 3
    row = i // 3
    add_tech_badge(sl, Inches(6.5) + col * Inches(1.8), Inches(2.85) + row * Inches(0.58), icon, name)

# Build tools section
add_text(sl, Inches(0.8), Inches(4.4), Inches(3), Inches(0.35),
         "⚙  BUILD & CONFIG", font_size=11, color=ACCENT_LAV, bold=True)

build_badges = [("📦", "Maven Wrapper"), ("🔧", "Spring Profiles"), ("📋", "JPA / Hibernate")]

for i, (icon, name) in enumerate(build_badges):
    add_tech_badge(sl, Inches(0.8) + i * Inches(1.8), Inches(4.85), icon, name)

# Architecture summary in a card
add_rect(sl, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.4), BG_CARD, CARD_BORDER, Pt(1))
add_text(sl, Inches(1.2), Inches(5.75), Inches(11), Inches(0.3),
         "Architecture Highlights", font_size=12, color=WHITE, bold=True)
add_text(sl, Inches(1.2), Inches(6.1), Inches(11), Inches(0.8),
         "• MVC pattern with layered separation: Controllers → Services → Repositories → Database\n"
         "• 15 REST Controllers  |  14 JPA Entities  |  DTO pattern for clean data transfer\n"
         "• Environment-based config with local-db.properties for flexible deployment",
         font_size=10, color=SUBTLE, line_spacing=1.7)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(5), Inches(-1), Inches(4), ACCENT_PURPLE, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Architecture")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "System Architecture", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "Clean layered architecture following MVC pattern with RESTful APIs.",
         font_size=12, color=SUBTLE)

# Architecture flow
flow_y = Inches(2.6)
boxes = [
    ("🖥️", "Browser", "HTML / CSS / JS"),
    ("🔁", "REST API", "JSON over HTTP"),
    ("🍃", "Spring Boot", "Controllers · Services"),
    ("📦", "JPA / Hibernate", "ORM Layer"),
    ("🐬", "MySQL 8", "Persistent Storage"),
]
x_start = Inches(0.6)
for i, (icon, title, sub) in enumerate(boxes):
    add_arch_box(sl, x_start + i * Inches(2.3), flow_y, icon, title, sub)
    if i < len(boxes) - 1:
        add_arrow(sl, x_start + Inches(1.7) + i * Inches(2.3), flow_y + Inches(0.15))

# Backend layers
add_text(sl, Inches(0.8), Inches(4.0), Inches(4), Inches(0.35),
         "Backend Components", font_size=12, color=WHITE, bold=True)

pills_data = [
    (ACCENT_PURPLE, "Controllers (15)"), (ACCENT_BLUE, "Entities (14)"),
    (ACCENT_CYAN, "Repositories"), (ACCENT_LAV, "Services"),
    (RGBColor(0xF5, 0x9E, 0x0B), "DTOs"), (RGBColor(0xEF, 0x44, 0x44), "Config"),
]

for i, (color, text) in enumerate(pills_data):
    col = i % 3
    row = i // 3
    add_pill(sl, Inches(0.8) + col * Inches(2.0), Inches(4.5) + row * Inches(0.48), color, text)

# Frontend pages
add_text(sl, Inches(7), Inches(4.0), Inches(4), Inches(0.35),
         "Frontend Pages (30+)", font_size=12, color=WHITE, bold=True)

page_groups = [
    (ACCENT_PURPLE, "Landing · Auth · Dashboard"),
    (ACCENT_BLUE, "Admissions · Faculty · Courses"),
    (ACCENT_CYAN, "Attendance (4 pages)"),
    (ACCENT_LAV, "Exams (4 pages)"),
    (RGBColor(0xF5, 0x9E, 0x0B), "Finance (4 pages)"),
    (RGBColor(0x10, 0xB9, 0x81), "Library (3) · Reports (4)"),
]

for i, (color, text) in enumerate(page_groups):
    col = i % 2
    row = i // 2
    add_pill(sl, Inches(7) + col * Inches(2.8), Inches(4.5) + row * Inches(0.48), color, text)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — CORE MODULES
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(9), Inches(0), Inches(4), ACCENT_BLUE, 0.05)
add_decorative_orb(sl, Inches(-1), Inches(5), Inches(3), ACCENT_LAV, 0.04)

add_tag(sl, Inches(0.8), Inches(0.5), "Features")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Core Modules", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "8 integrated modules covering every university operation.",
         font_size=12, color=SUBTLE)

modules = [
    ("🎓", "Admissions", "Student Registration"),
    ("👨‍🏫", "Faculty", "Staff Management"),
    ("📚", "Courses", "Curriculum Setup"),
    ("🏛️", "Departments", "Org Structure"),
    ("📋", "Attendance", "Student & Teacher"),
    ("📝", "Examinations", "Schedule · Marks · Grades"),
    ("💰", "Finance", "Fees · Scholarships"),
    ("📖", "Library", "Books · Issues · Members"),
]

for i, (icon, name, sub) in enumerate(modules):
    col = i % 4
    row = i // 4
    add_module_card(sl, Inches(0.8) + col * Inches(3.0), Inches(2.4) + row * Inches(1.4),
                    Inches(2.7), Inches(1.15), icon, name, sub)

# Stats row at bottom
stats_y = Inches(5.4)
add_rect(sl, Inches(0.8), stats_y, Inches(11.7), Inches(1.5), BG_CARD, CARD_BORDER, Pt(1))

stat_items = [("30+", "Pages"), ("15", "Controllers"), ("14", "Entities"), ("10+", "API Routes"), ("8", "Modules")]
for i, (num, label) in enumerate(stat_items):
    add_stat_box(sl, Inches(1.0) + i * Inches(2.3), stats_y + Inches(0.25), num, label)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — DATABASE DESIGN
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(10), Inches(3), Inches(4), ACCENT_PURPLE, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Data Model")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Database Design", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "14 JPA entities mapped to MySQL 8 with relational integrity.",
         font_size=12, color=SUBTLE)

entities = [
    ("👤", "User", "Auth, roles, passwords"),
    ("🎓", "Student", "Enrollment, department"),
    ("👨‍🏫", "Faculty", "Designation, qualifications"),
    ("📘", "Course", "Credits, semester mapping"),
    ("🏛️", "Department", "Name, HOD, date"),
    ("📋", "Attendance", "Daily records"),
    ("📅", "ExamSchedule", "Date, venue, course"),
    ("📊", "ExamResult", "Marks, grades, CGPA"),
    ("💳", "FeePayment", "Amount, date, status"),
    ("📄", "FeeStructure", "Dept-wise breakdown"),
    ("🏆", "Scholarship", "Criteria, eligibility"),
    ("📖", "LibraryBook", "Catalog management"),
    ("🔄", "LibraryIssue", "Issue / return tracking"),
    ("📢", "Announcement", "System notifications"),
]

for i, (icon, name, desc) in enumerate(entities):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.1)
    y = Inches(2.4) + row * Inches(1.5)
    card_w, card_h = Inches(2.85), Inches(1.2)
    
    card = add_rect(sl, x, y, card_w, card_h, BG_CARD, CARD_BORDER, Pt(1))
    add_text(sl, x + Inches(0.15), y + Inches(0.12), Inches(0.4), Inches(0.35),
             icon, font_size=16, align=PP_ALIGN.LEFT)
    add_text(sl, x + Inches(0.5), y + Inches(0.12), card_w - Inches(0.6), Inches(0.3),
             name, font_size=11, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.5), y + Inches(0.45), card_w - Inches(0.6), Inches(0.5),
             desc, font_size=8.5, color=DIM)

# last row has only 2 items, so add a note card
add_rect(sl, Inches(6.7), Inches(5.4), Inches(5.9), Inches(1.2), BG_CARD, CARD_BORDER, Pt(1))
add_text(sl, Inches(7.0), Inches(5.55), Inches(5.3), Inches(0.3),
         "✦  Relational Mapping", font_size=11, color=ACCENT_CYAN, bold=True)
add_text(sl, Inches(7.0), Inches(5.9), Inches(5.3), Inches(0.6),
         "All entities use JPA annotations (@Entity, @ManyToOne, @OneToMany)\nfor automatic schema generation and relational integrity.",
         font_size=9, color=SUBTLE, line_spacing=1.5)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — REST API ENDPOINTS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(-1), Inches(-1), Inches(4), ACCENT_BLUE, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Backend")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "REST API Endpoints", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "Clean RESTful API design with consistent JSON responses.",
         font_size=12, color=SUBTLE)

# Left column
left_apis = [
    (ACCENT_PURPLE, "/api/students — CRUD operations"),
    (ACCENT_BLUE,   "/api/faculty — Staff management"),
    (ACCENT_CYAN,   "/api/departments — Org units"),
    (ACCENT_LAV,    "/api/courses — Curriculum"),
    (RGBColor(0xF5, 0x9E, 0x0B), "/api/attendance — Daily records"),
]

y_pos = Inches(2.4)
for dot_color, text in left_apis:
    dy = add_feature_row(sl, Inches(0.8), y_pos, dot_color, text, w=Inches(5.4))
    y_pos += dy

# Right column
right_apis = [
    (ACCENT_PURPLE, "/api/exams — Schedule management"),
    (ACCENT_BLUE,   "/api/exam-results — Grading system"),
    (ACCENT_CYAN,   "/api/fees — Payment processing"),
    (ACCENT_LAV,    "/api/scholarships — Award management"),
    (RGBColor(0xF5, 0x9E, 0x0B), "/api/library/* — Books & Issues"),
]

y_pos = Inches(2.4)
for dot_color, text in right_apis:
    dy = add_feature_row(sl, Inches(6.8), y_pos, dot_color, text, w=Inches(5.4))
    y_pos += dy

# API info card
add_rect(sl, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.7), BG_CARD, CARD_BORDER, Pt(1))
add_text(sl, Inches(1.2), Inches(5.4), Inches(4), Inches(0.3),
         "API Design Principles", font_size=12, color=WHITE, bold=True)

principles = [
    "✓  RESTful conventions — GET, POST, PUT, DELETE",
    "✓  JSON request/response bodies with proper HTTP status codes",
    "✓  Session-based authentication with secure cookie management",
    "✓  Cross-origin support for frontend-backend communication",
]
for i, p in enumerate(principles):
    add_text(sl, Inches(1.2), Inches(5.8) + i * Inches(0.27), Inches(10), Inches(0.25),
             p, font_size=9.5, color=SUBTLE)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — UI/UX DESIGN
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(8), Inches(-1), Inches(5), ACCENT_PURPLE, 0.07)
add_decorative_orb(sl, Inches(-1), Inches(4), Inches(4), ACCENT_CYAN, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Design")

add_text(sl, Inches(0.8), Inches(1.0), Inches(10), Inches(0.7),
         "UI / UX Philosophy", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "A premium, immersive interface that makes university management delightful.",
         font_size=12, color=SUBTLE)

ui_cards = [
    ("🌌", "Glassmorphism", "Frosted-glass cards with subtle transparency, blur effects, and glowing borders for visual depth."),
    ("✨", "Particle Effects", "Interactive particle canvas across landing, auth, and dashboard pages for a living feel."),
    ("🎨", "Gradient System", "Curated purple-blue-cyan palette with smooth gradients that unify the entire experience."),
    ("💫", "Micro-Animations", "Hover effects, page transitions, logo animations, and toast notifications for polish."),
]

accents = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_LAV]
for i, (icon, title, desc) in enumerate(ui_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.8)
    y = Inches(2.4) + row * Inches(2.2)
    add_glass_card(sl, x, y, Inches(5.4), Inches(1.9), icon, title, desc)
    # accent bar
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.4), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accents[i]
    bar.line.fill.background()

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — USER FLOW
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(10), Inches(1), Inches(3), ACCENT_LAV, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Workflow")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "User Flow", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "Seamless navigation from landing to any management module.",
         font_size=12, color=SUBTLE)

# Timeline
steps = [
    ("1", "Landing Page", "Animated logo with particle canvas → \"Enter the System\" button with glow effects"),
    ("2", "Authentication", "Glassmorphic login card → credential validation → \"Access Granted\" overlay"),
    ("3", "Dashboard", "Stats grid (students, faculty, attendance, fees) → quick actions → announcements"),
    ("4", "Module Pages", "Sidebar navigation → CRUD forms → data tables → search & filter → toast alerts"),
    ("5", "Reports", "Student, attendance, finance & exam reports with aggregated analytics"),
]

# Draw the timeline vertical line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(2.4), Inches(0.04), Inches(4.7))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_PURPLE
line.line.fill.background()

for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.4) + i * Inches(0.95)
    
    # Dot on timeline
    dot = add_circle(sl, Inches(1.27), y + Inches(0.12), Inches(0.3), ACCENT_PURPLE)
    add_text(sl, Inches(1.27), y + Inches(0.12), Inches(0.3), Inches(0.3),
             num, font_size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Content card
    card = add_rect(sl, Inches(1.9), y, Inches(9.8), Inches(0.78), BG_CARD, CARD_BORDER, Pt(0.75))
    add_text(sl, Inches(2.15), y + Inches(0.08), Inches(3), Inches(0.3),
             title, font_size=12, color=WHITE, bold=True)
    add_text(sl, Inches(2.15), y + Inches(0.38), Inches(9.2), Inches(0.35),
             desc, font_size=9.5, color=SUBTLE)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — KEY HIGHLIGHTS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(-1), Inches(2), Inches(4), ACCENT_CYAN, 0.06)

add_tag(sl, Inches(0.8), Inches(0.5), "Highlights")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Key Highlights", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "What sets UniVerse apart from conventional management systems.",
         font_size=12, color=SUBTLE)

# Left features
highlights = [
    (ACCENT_PURPLE, "30+ pages — comprehensive coverage of all operations"),
    (ACCENT_BLUE,   "RESTful APIs — clean, documented JSON endpoints"),
    (ACCENT_CYAN,   "Real-time Dashboard — live stats with animated counters"),
    (ACCENT_LAV,    "Secure Auth — password encryption via Spring Security"),
    (RGBColor(0x10, 0xB9, 0x81), "Library System — full book catalog & issue tracking"),
    (RGBColor(0xF5, 0x9E, 0x0B), "Finance Module — fees, scholarships & reporting"),
]

y_pos = Inches(2.4)
for dot_color, text in highlights:
    dy = add_feature_row(sl, Inches(0.8), y_pos, dot_color, text, w=Inches(6.5))
    y_pos += dy

# Right stats
stats = [("30+", "Pages"), ("15", "Controllers"), ("14", "Entities"), ("8", "Modules")]
for i, (num, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = Inches(8.2) + col * Inches(2.3)
    y = Inches(2.5) + row * Inches(1.8)
    add_rect(sl, x, y, Inches(2.1), Inches(1.4), BG_CARD, CARD_BORDER, Pt(1))
    add_stat_box(sl, x + Inches(0.15), y + Inches(0.2), num, label)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — FUTURE SCOPE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(8), Inches(4), Inches(4), ACCENT_PURPLE, 0.06)

add_tag(sl, Inches(0.8), Inches(0.5), "Roadmap")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Future Scope", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(7), Inches(0.4),
         "Planned enhancements to evolve UniVerse into a production-ready platform.",
         font_size=12, color=SUBTLE)

future_cards = [
    ("📱", "Mobile Responsive", "Progressive Web App with native-like experience on phones and tablets."),
    ("🔔", "Notifications", "Email & push notifications for fee reminders, exam updates & announcements."),
    ("📈", "Advanced Analytics", "Interactive charts, performance trends, and predictive insights with Chart.js."),
    ("👥", "Multi-Role Access", "Student, teacher, and parent portals with role-specific dashboards."),
]

future_accents = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_LAV]
for i, (icon, title, desc) in enumerate(future_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.8)
    y = Inches(2.4) + row * Inches(2.2)
    add_glass_card(sl, x, y, Inches(5.4), Inches(1.9), icon, title, desc)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.4), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = future_accents[i]
    bar.line.fill.background()

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)
add_decorative_orb(sl, Inches(-1), Inches(-1), Inches(5), ACCENT_PURPLE, 0.07)
add_decorative_orb(sl, Inches(9), Inches(4), Inches(4), ACCENT_CYAN, 0.05)

add_tag(sl, Inches(0.8), Inches(0.5), "Summary")

add_text(sl, Inches(0.8), Inches(1.0), Inches(8), Inches(0.7),
         "Conclusion", font_size=36, color=ACCENT_PURPLE, bold=True)

add_text(sl, Inches(0.8), Inches(1.7), Inches(10), Inches(0.8),
         "UniVerse demonstrates that university management can be beautiful, efficient, and unified. "
         "By combining a robust Spring Boot backend with an immersive glassmorphism frontend, "
         "we've built a system that simplifies operations while delighting users.",
         font_size=12, color=SUBTLE, line_spacing=1.7)

# 3 conclusion cards
concl_cards = [
    ("🎯", "Problem Solved", "Unified fragmented university operations into one integrated platform"),
    ("⚙️", "Full-Stack", "Java 17 + Spring Boot + MySQL + Modern Frontend with particle effects"),
    ("🚀", "Production Ready", "Scalable MVC architecture with room for growth and enhancement"),
]

for i, (icon, title, desc) in enumerate(concl_cards):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(3.0)
    card = add_rect(sl, x, y, Inches(3.7), Inches(2.0), BG_CARD, CARD_BORDER, Pt(1))
    
    # accent top bar
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.7), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN][i]
    bar.line.fill.background()
    
    add_text(sl, x, y + Inches(0.25), Inches(3.7), Inches(0.45),
             icon, font_size=28, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.2), y + Inches(0.85), Inches(3.3), Inches(0.3),
             title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.2), y + Inches(1.2), Inches(3.3), Inches(0.6),
             desc, font_size=9.5, color=SUBTLE, align=PP_ALIGN.CENTER, line_spacing=1.5)

# Key takeaway bar
add_rect(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), BG_CARD, CARD_BORDER, Pt(1))
add_gradient_bar(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.035))
add_text(sl, Inches(1.2), Inches(5.7), Inches(11), Inches(0.3),
         "Key Takeaway", font_size=12, color=WHITE, bold=True)
add_text(sl, Inches(1.2), Inches(6.05), Inches(11), Inches(0.5),
         "UniVerse is not just a project — it's a vision for how educational institutions can leverage "
         "modern technology to streamline their operations and provide a world-class digital experience.",
         font_size=10, color=SUBTLE, line_spacing=1.6)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(sl, BG_DARK)

# Big decorative orbs
add_decorative_orb(sl, Inches(0), Inches(-1), Inches(6), ACCENT_PURPLE, 0.08)
add_decorative_orb(sl, Inches(8), Inches(3), Inches(6), ACCENT_CYAN, 0.06)
add_decorative_orb(sl, Inches(3), Inches(5), Inches(4), ACCENT_BLUE, 0.04)

# Logo
if os.path.exists(logo_path):
    # Logo background circle
    logo_bg = add_circle(sl, Inches(5.9), Inches(0.4), Inches(1.4), RGBColor(0x14, 0x14, 0x35))
    logo_bg.line.fill.solid()
    logo_bg.line.fill.fore_color.rgb = RGBColor(0x3A, 0x35, 0x70)
    logo_bg.line.width = Pt(2)
    sl.shapes.add_picture(logo_path, Inches(6.15), Inches(0.65), Inches(0.9), Inches(0.9))

# Thank You text
add_text(sl, Inches(0), Inches(2.0), SLIDE_W, Inches(1.0),
         "Thank You!", font_size=52, color=ACCENT_PURPLE, bold=True,
         align=PP_ALIGN.CENTER, font_name='Calibri')

add_text(sl, Inches(0), Inches(2.9), SLIDE_W, Inches(0.4),
         "We appreciate your time and attention",
         font_size=14, color=SUBTLE, align=PP_ALIGN.CENTER)

add_gradient_bar(sl, Inches(5.2), Inches(3.5), Inches(2.9), Inches(0.03))

# ── Team Cards ──────────────────────────────────────────────────
card_w, card_h = Inches(4.2), Inches(2.4)

# Card 1 — Prachi
c1_x = Inches(2.3)
c1_y = Inches(4.0)
card1 = add_rect(sl, c1_x, c1_y, card_w, card_h, BG_CARD, CARD_BORDER, Pt(1))
# Top accent
bar1 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, c1_x, c1_y, card_w, Inches(0.04))
bar1.fill.solid()
bar1.fill.fore_color.rgb = ACCENT_PURPLE
bar1.line.fill.background()

# Avatar circle
av1 = add_circle(sl, c1_x + Inches(1.5), c1_y + Inches(0.25), Inches(1.0), ACCENT_PURPLE)
add_text(sl, c1_x + Inches(1.5), c1_y + Inches(0.4), Inches(1.0), Inches(0.7),
         "PP", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri')

# Name
add_text(sl, c1_x, c1_y + Inches(1.35), card_w, Inches(0.35),
         "Prachi Pandey", font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ID
add_text(sl, c1_x, c1_y + Inches(1.72), card_w, Inches(0.3),
         "S25CSEU1097", font_size=11, color=ACCENT_CYAN, align=PP_ALIGN.CENTER, bold=True)

# Role
add_text(sl, c1_x, c1_y + Inches(2.0), card_w, Inches(0.3),
         "Full-Stack Developer", font_size=9, color=DIM, align=PP_ALIGN.CENTER)

# Card 2 — Anushka
c2_x = Inches(6.8)
c2_y = Inches(4.0)
card2 = add_rect(sl, c2_x, c2_y, card_w, card_h, BG_CARD, CARD_BORDER, Pt(1))
# Top accent
bar2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, c2_x, c2_y, card_w, Inches(0.04))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_CYAN
bar2.line.fill.background()

# Avatar circle
av2 = add_circle(sl, c2_x + Inches(1.5), c2_y + Inches(0.25), Inches(1.0), ACCENT_BLUE)
add_text(sl, c2_x + Inches(1.5), c2_y + Inches(0.4), Inches(1.0), Inches(0.7),
         "A", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri')

# Name
add_text(sl, c2_x, c2_y + Inches(1.35), card_w, Inches(0.35),
         "Anushka", font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ID
add_text(sl, c2_x, c2_y + Inches(1.72), card_w, Inches(0.3),
         "S25CSEU1113", font_size=11, color=ACCENT_CYAN, align=PP_ALIGN.CENTER, bold=True)

# Role
add_text(sl, c2_x, c2_y + Inches(2.0), card_w, Inches(0.3),
         "Full-Stack Developer", font_size=9, color=DIM, align=PP_ALIGN.CENTER)

# Footer
add_text(sl, Inches(0), Inches(6.7), SLIDE_W, Inches(0.3),
         "UniVerse — University Management System  ·  April 2026",
         font_size=9, color=DIM, align=PP_ALIGN.CENTER)

add_gradient_bar(sl, Inches(0), Inches(7.42), SLIDE_W, Inches(0.05))


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'UniVerse_Presentation.pptx')
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: 13.333\" × 7.5\" (Widescreen 16:9)")
